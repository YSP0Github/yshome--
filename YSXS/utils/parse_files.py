import re
from datetime import datetime
from typing import Dict, List, Optional, Sequence

from PyPDF2 import PdfReader
from docx import Document as DocxDocument

from .datetimes import format_cn_time, to_cst, utc_now
MAX_TEXT_CHARS = 8000
CONTROL_CHARS_RE = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]")
YEAR_PATTERN = re.compile(r"((?:18|19|20)\d{2})")
KEYWORD_SPLIT_PATTERN = re.compile(r"[;,，、；]+")
ABSTRACT_PATTERNS = [
    r"(?:abstract|summary)\s*[:：]\s*(.+)",
    r"摘要\s*[:：]\s*(.+)",
    r"摘要\s*(.+)",
]
KEYWORD_PATTERNS = [
    r"(?:keywords?|key\s*words?)\s*[:：]\s*([^\n]+)",
    r"关键词\s*[:：]\s*([^\n]+)",
    r"关键词\s*(.+)",
]
JOURNAL_PATTERNS = [
    r"(?:journal|期刊)\s*[:：]\s*([^\n，。,;]+)",
    r"发表于\s*(?:《)?([^\s，。]+)",
    r"published in\s+([^\n,.;]+)",
    r"proceedings of\s+([^\n,.;]+)",
    r"(?:会议|conference)\s*[:：]\s*([^\n，。,;]+)",
]
AUTHOR_PATTERNS = [
    r"作者[:：]\s*([^\n]+)",
    r"作者名单[:：]\s*([^\n]+)",
    r"authors?\s*[:：]\s*([^\n]+)",
    r"by\s+([^\n]+)",
]
DOC_TYPE_RULES = [
    ("conference", "conference|symposium|proceedings|workshop|研讨会|会议|论坛"),
    ("review", "review|reviews|review article|literature review|annual review|survey|overview|state of the art|progress|systematic review|meta-analysis|meta analysis|综述|评述|研究现状|文献综述|综述论文"),
    ("preprint", "preprint|working paper|arxiv|预印本"),
    ("patent", "patent|专利"),
    ("standard", "standard|specification|标准"),
    ("dataset", "dataset|data set|database|数据集"),
    ("software", "software|code|computer program|toolbox|package|软件"),
    ("thesis", "thesis|dissertation|学位论文|博士论文|硕士论文"),
    ("report", "report|white paper|年度报告|报告"),
    ("book", "monograph|手册|manual|book|教材|专著"),
]
CATEGORY_KEYWORDS = {
    "ai": [
        "artificial intelligence",
        "machine learning",
        "deep learning",
        "neural network",
        "人工智能",
        "机器学习",
        "深度学习",
    ],
    "computer": [
        "computer",
        "algorithm",
        "software",
        "计算机",
        "算法",
        "程序设计",
        "网络",
    ],
    "math": [
        "mathematics",
        "math",
        "建模",
        "模型",
        "数值分析",
        "statistics",
    ],
    "physics": ["physics", "物理", "力学", "quantum"],
    "geology": ["geology", "地质", "岩石", "矿物"],
    "geophysics": ["geophysics", "地球物理", "地震", "勘探"],
    "geochemistry": ["geochemistry", "地球化学", "同位素"],
    "geography": ["geography", "地理", "遥感"],
    "geoinformatics": ["gis", "地理信息", "遥感技术", "geoinformatics"],
    "geostatistics": ["geostatistics", "地质统计"],
    "biology": ["biology", "生物", "基因", "生命科学"],
}
ALLOWED_TEXT_ENCODINGS = ["utf-8", "utf-16", "gbk", "big5", "latin-1"]


def default_result() -> Dict[str, str]:
    return {
        "title": "",
        "authors": "",
        "journal": "",
        "year": "",
        "keywords": "",
        "abstract": "",
        "full_text": "",
        "category": "",
        "type": "",
        "error": "",
    }


def limit_length(value: str, limit: int = 255) -> str:
    if not value:
        return ""
    value = value.strip()
    return value if len(value) <= limit else value[: limit - 1].rstrip() + "…"


def clean_text(text: str) -> str:
    if not text:
        return ""
    normalized = text.replace("\r\n", "\n").replace("\r", "\n")
    normalized = CONTROL_CHARS_RE.sub("", normalized)
    normalized = re.sub(r"[ \t]+", " ", normalized)
    normalized = re.sub(r"\n{3,}", "\n\n", normalized)
    return normalized.strip()


def normalize_keywords(raw: str) -> str:
    if not raw:
        return ""
    parts = [segment.strip(" ：:，,;") for segment in KEYWORD_SPLIT_PATTERN.split(raw)]
    parts = [p for p in parts if p]
    seen: List[str] = []
    for keyword in parts:
        if keyword.lower() not in [s.lower() for s in seen]:
            seen.append(keyword)
    return limit_length(", ".join(seen))


def _match_pattern_list(text: str, patterns: Sequence[str], flags=re.IGNORECASE) -> str:
    if not text:
        return ""
    for pattern in patterns:
        match = re.search(pattern, text, flags)
        if match:
            candidate = match.group(1).strip()
            if candidate:
                return candidate
    return ""


def extract_title(candidate_texts: Sequence[str]) -> str:
    filtered = [t.strip() for t in candidate_texts if t and 5 <= len(t.strip()) <= 180]
    if not filtered:
        filtered = [t.strip() for t in candidate_texts if t]
    for text in filtered[:5]:
        upper_ratio = sum(1 for c in text if c.isupper()) / max(len(text), 1)
        if upper_ratio > 0.3:
            return limit_length(text, 200)
    return limit_length(filtered[0], 200) if filtered else ""


def extract_authors(text: str, fallback_lines: Optional[Sequence[str]] = None) -> str:
    if not text:
        return ""
    matched = _match_pattern_list(text, AUTHOR_PATTERNS)
    if matched:
        return limit_length(_normalize_author_list(matched))
    candidate_zone = " ".join(text.split("\n")[:4])
    if "," in candidate_zone or "，" in candidate_zone:
        return limit_length(_normalize_author_list(candidate_zone.split(".", 1)[0]))

    if fallback_lines:
        for line in fallback_lines:
            plain = re.sub(r"[\d\(\)\[\]·•]", "", line).strip()
            if not plain or len(plain) > 160:
                continue
            if any(keyword in plain.lower() for keyword in ("abstract", "summary", "keywords")):
                continue
            if any(keyword in plain for keyword in ("摘要", "关键词", "作者", "单位")):
                continue
            separators = [",", ";", "，", "、", "；"]
            if sum(sep in plain for sep in separators) >= 1:
                return limit_length(_normalize_author_list(plain))
    return ""


def _normalize_author_list(raw: str) -> str:
    tokens = re.split(r"\s*(?:,|;|，|、|；| and |&|与)\s*", raw)
    tokens = [token.strip() for token in tokens if token.strip()]
    seen: List[str] = []
    for token in tokens:
        normalized = token.strip(" .")
        if not normalized:
            continue
        if normalized.lower() not in [s.lower() for s in seen]:
            seen.append(normalized)
    return ", ".join(seen)


def extract_keywords(text: str) -> str:
    matched = _match_pattern_list(text, KEYWORD_PATTERNS)
    return normalize_keywords(matched)


def extract_journal(text: str) -> str:
    matched = _match_pattern_list(text, JOURNAL_PATTERNS)
    return limit_length(matched)


def extract_year(*candidates: str) -> str:
    current_year = to_cst(utc_now()).year + 1
    for candidate in candidates:
        if not candidate:
            continue
        match = YEAR_PATTERN.search(str(candidate))
        if match:
            year_value = int(match.group(1))
            if 1800 <= year_value <= current_year:
                return str(year_value)
    return ""


def extract_abstract(text: str) -> str:
    if not text:
        return ""
    for pattern in ABSTRACT_PATTERNS:
        match = re.search(pattern, text, re.IGNORECASE | re.DOTALL)
        if match:
            abstract = match.group(1).strip()
            if "\n\n" in abstract:
                abstract = abstract.split("\n\n", 1)[0]
            return limit_length(abstract, 1200)
    snippet = text[:1200]
    return limit_length(snippet, 1200)


def guess_doc_type(full_text: str, journal_line: str) -> str:
    haystack = " ".join(filter(None, [journal_line, full_text])).lower()
    for doc_type, pattern in DOC_TYPE_RULES:
        if re.search(pattern, haystack, re.IGNORECASE):
            return doc_type
    if "journal" in haystack or "期刊" in haystack:
        return "journal"
    return "other"


def guess_category_from_text(*parts: str) -> str:
    haystack = " ".join(filter(None, parts))
    normalized = haystack.lower()
    for slug, keywords in CATEGORY_KEYWORDS.items():
        for keyword in keywords:
            if keyword.lower() in normalized:
                return slug
    return ""


def combine_sections(sections: Sequence[str]) -> str:
    combined = "\n".join([section for section in sections if section])
    if len(combined) > MAX_TEXT_CHARS:
        return combined[:MAX_TEXT_CHARS]
    return combined


def build_line_candidates(blocks: Sequence[str], limit: int = 8) -> List[str]:
    lines: List[str] = []
    for block in blocks:
        if not block:
            continue
        for raw_line in block.split("\n"):
            line = raw_line.strip()
            if 5 <= len(line) <= 160:
                lines.append(line)
            if len(lines) >= limit:
                return lines
    return lines


def parse_pdf(file_path: str) -> Dict[str, str]:
    result = default_result()
    try:
        reader = PdfReader(file_path)
        metadata = reader.metadata or {}
        metadata_title = clean_text(str(metadata.get("/Title", "")))
        metadata_author = clean_text(str(metadata.get("/Author", "")))
        metadata_keywords = clean_text(str(metadata.get("/Keywords", "")))
        pages_text: List[str] = []
        for page in reader.pages[:8]:
            try:
                page_text = page.extract_text() or ""
            except Exception:
                continue
            cleaned = clean_text(page_text)
            if cleaned:
                pages_text.append(cleaned)
        full_text = combine_sections(pages_text)
        line_candidates = build_line_candidates(pages_text[:2], limit=10)
        title = extract_title([metadata_title] + line_candidates)
        authors = extract_authors(
            f"{metadata_author}\n{full_text[:500]}",
            fallback_lines=line_candidates[1:5],
        )
        keywords = normalize_keywords(metadata_keywords or extract_keywords(full_text))
        journal = extract_journal(full_text)
        year = extract_year(metadata.get("/CreationDate"), metadata.get("/ModDate"), full_text)
        abstract = extract_abstract(full_text)
        doc_type = guess_doc_type(full_text, journal)
        category = guess_category_from_text(keywords, journal, full_text)
        result.update(
            {
                "title": title,
                "authors": authors,
                "journal": journal,
                "year": year,
                "keywords": keywords,
                "abstract": abstract,
                "full_text": full_text,
                "type": doc_type,
                "category": category,
            }
        )
    except Exception as exc:
        result["error"] = f"PDF解析错误: {exc}"
    return result


def parse_docx(file_path: str) -> Dict[str, str]:
    result = default_result()
    try:
        doc = DocxDocument(file_path)
        paragraphs = [clean_text(p.text) for p in doc.paragraphs if clean_text(p.text)]
        full_text = combine_sections(paragraphs)
        core = doc.core_properties
        meta_title = clean_text(core.title or "")
        meta_author = clean_text(core.author or "")
        meta_keywords = clean_text(core.keywords or "")
        line_candidates = build_line_candidates(paragraphs[:6], limit=10)
        title = extract_title([meta_title] + line_candidates)
        authors = extract_authors(
            f"{meta_author}\n{full_text[:400]}",
            fallback_lines=line_candidates[1:5],
        )
        keywords = normalize_keywords(meta_keywords or extract_keywords(full_text))
        journal = extract_journal(full_text)
        year = extract_year(core.created, core.last_printed, full_text)
        abstract = extract_abstract(full_text)
        doc_type = guess_doc_type(full_text, journal)
        category = guess_category_from_text(keywords, journal, full_text)
        result.update(
            {
                "title": title,
                "authors": authors,
                "journal": journal,
                "year": year,
                "keywords": keywords,
                "abstract": abstract,
                "full_text": full_text,
                "type": doc_type,
                "category": category,
            }
        )
    except Exception as exc:
        result["error"] = f"DOCX解析错误: {exc}"
    return result


def parse_pptx(file_path: str) -> Dict[str, str]:
    result = default_result()
    try:
        from pptx import Presentation

        presentation = Presentation(file_path)
        slides_text: List[str] = []
        for slide in presentation.slides[:5]:
            shape_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    cleaned = clean_text(shape.text)
                    if cleaned:
                        shape_text.append(cleaned)
            if shape_text:
                slides_text.append("\n".join(shape_text))
        full_text = combine_sections(slides_text)
        core = presentation.core_properties
        meta_title = clean_text(core.title or "")
        meta_author = clean_text(core.author or "")
        meta_keywords = clean_text(core.keywords or "")
        line_candidates = build_line_candidates(slides_text[:3], limit=10)
        title = extract_title([meta_title] + line_candidates)
        authors = extract_authors(
            f"{meta_author}\n{full_text[:400]}",
            fallback_lines=line_candidates[1:5],
        )
        keywords = normalize_keywords(meta_keywords or extract_keywords(full_text))
        journal = extract_journal(full_text)
        year = extract_year(core.created, core.last_printed, full_text)
        abstract = extract_abstract(full_text)
        doc_type = guess_doc_type(full_text, journal)
        category = guess_category_from_text(keywords, journal, full_text)
        result.update(
            {
                "title": title,
                "authors": authors,
                "journal": journal,
                "year": year,
                "keywords": keywords,
                "abstract": abstract,
                "full_text": full_text,
                "type": doc_type,
                "category": category,
            }
        )
    except Exception as exc:
        result["error"] = f"PPTX解析错误: {exc}"
    return result


def parse_txt(file_path: str) -> Dict[str, str]:
    result = default_result()
    content = ""
    try:
        for encoding in ALLOWED_TEXT_ENCODINGS:
            try:
                with open(file_path, "r", encoding=encoding) as handle:
                    content = handle.read()
                break
            except UnicodeDecodeError:
                continue
        if not content:
            raise UnicodeDecodeError("unknown", b"", 0, 1, "无法解码文本文件")
        cleaned = clean_text(content)
        paragraphs = [segment for segment in cleaned.split("\n") if segment.strip()]
        full_text = combine_sections(paragraphs)
        line_candidates = build_line_candidates(paragraphs[:8], limit=10)
        title = extract_title(line_candidates or paragraphs[:3])
        authors = extract_authors(
            full_text[:400],
            fallback_lines=line_candidates[1:5],
        )
        keywords = normalize_keywords(extract_keywords(full_text))
        journal = extract_journal(full_text)
        year = extract_year(full_text)
        abstract = extract_abstract(full_text)
        doc_type = guess_doc_type(full_text, journal)
        category = guess_category_from_text(keywords, journal, full_text)
        result.update(
            {
                "title": title,
                "authors": authors,
                "journal": journal,
                "year": year,
                "keywords": keywords,
                "abstract": abstract,
                "full_text": full_text,
                "type": doc_type,
                "category": category,
            }
        )
    except Exception as exc:
        result["error"] = f"TXT解析错误: {exc}"
    return result


def parse_files(file_path: str, file_type: str) -> Dict[str, str]:
    parser_map = {
        "pdf": parse_pdf,
        "doc": parse_txt,
        "docx": parse_docx,
        "ppt": parse_pptx,
        "pptx": parse_pptx,
        "txt": parse_txt,
    }
    result = default_result()
    parser = parser_map.get((file_type or "").lower())
    if not parser:
        result["error"] = f"不支持的文件类型: {file_type}"
        return result
    parsed = parser(file_path)
    normalized = default_result()
    normalized.update(parsed)
    normalized["year"] = str(normalized["year"]).strip() if normalized["year"] else ""
    normalized["parsed_at"] = format_cn_time(utc_now(), "%Y-%m-%d %H:%M:%S")
    normalized["keywords"] = limit_length(normalized.get("keywords", ""))
    normalized["authors"] = limit_length(normalized.get("authors", ""))
    normalized["title"] = limit_length(normalized.get("title", ""), 200)
    return normalized
